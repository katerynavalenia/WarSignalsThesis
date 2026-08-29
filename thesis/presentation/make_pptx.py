"""Build ``defence.pptx`` from the content spine in ``slides.py``.

Run from anywhere::

    python thesis/presentation/make_pptx.py

Everything is drawn on blank layouts with explicit geometry rather than through
PowerPoint's placeholders, because placeholder inheritance is what makes a deck
look different on the machine it is presented from than on the one it was built
on. The cost is that positions are computed here; the benefit is that the file
opens the same way in PowerPoint, Keynote and Google Slides.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

import slides as content            # noqa: E402
from markup import to_runs          # noqa: E402

HERE = Path(__file__).resolve().parent
FIGURES = HERE / "figures"
OUTPUT = HERE / "defence.pptx"

# ---------------------------------------------------------------------------
# Palette and metrics
# ---------------------------------------------------------------------------

INK = RGBColor(0x16, 0x22, 0x2E)
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)       # the WAERLST blue of Figure 1
MUTED = RGBColor(0x5A, 0x6B, 0x7A)
RULE = RGBColor(0xCB, 0xD5, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF2, 0xF6, 0xF9)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
BODY_W = SLIDE_W - 2 * MARGIN

SECTION_Y = Inches(0.30)
TITLE_Y = Inches(0.56)
RULE_Y = Inches(1.36)
BODY_Y = Inches(1.56)
BODY_BOTTOM = Inches(7.00)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def _write(paragraph, text: str, size: int, color: RGBColor, bold=False, italic=False):
    """Emit the authored markup as runs, so bold survives inside a sentence."""
    for run in to_runs(text):
        r = paragraph.add_run()
        r.text = run.text
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = color
        r.font.bold = bold or run.bold
        r.font.italic = italic or run.italic


def _rule(slide, y, width=None, color=RULE, thickness=Pt(1.25)):
    line = slide.shapes.add_shape(1, MARGIN, y, width or BODY_W, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


# ---------------------------------------------------------------------------
# Slide chrome
# ---------------------------------------------------------------------------

def _chrome(slide, item: content.Slide, number: int, total: int):
    if item.section:
        frame = _textbox(slide, MARGIN, SECTION_Y, BODY_W, Inches(0.26))
        p = frame.paragraphs[0]
        _write(p, item.section.upper(), 11, ACCENT, bold=True)

    frame = _textbox(slide, MARGIN, TITLE_Y, BODY_W, Inches(0.72))
    p = frame.paragraphs[0]
    _write(p, item.title, 25, INK, bold=True)

    _rule(slide, RULE_Y)

    frame = _textbox(slide, SLIDE_W - MARGIN - Inches(1.2), Inches(7.05), Inches(1.2), Inches(0.25))
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _write(p, f"{number} / {total}", 10, MUTED)


def _bullet_size(bullets: list[content.Bullet]) -> int:
    """Shrink the body a step at a time rather than letting it overflow."""
    chars = sum(len(b.text) for b in bullets)
    if chars > 1500:
        return 12
    if chars > 1150:
        return 13
    if chars > 800:
        return 14
    return 16


def _bullets(slide, bullets: list[content.Bullet], top, height, size=None):
    size = size or _bullet_size(bullets)
    frame = _textbox(slide, MARGIN, top, BODY_W, height)
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_after = Pt(size * 0.62)
        p.line_spacing = 1.08
        indent = Inches(0.30) * bullet.level
        p.left_indent = Emu(int(indent) + int(Inches(0.20)))
        p.first_line_indent = -Inches(0.20)
        marker = "–  " if bullet.level else "▪  "
        _write(p, marker, size, ACCENT if not bullet.level else MUTED)
        _write(p, bullet.text, size - bullet.level, INK)
    return frame


def _figure(slide, name: str, caption: str, top, available):
    path = FIGURES / name
    if not path.exists():
        raise FileNotFoundError(path)
    pic = slide.shapes.add_picture(str(path), MARGIN, top, width=BODY_W)
    caption_h = Inches(0.34) if caption else Inches(0)
    if pic.height > available - caption_h:
        factor = (available - caption_h) / pic.height
        pic.width = int(pic.width * factor)
        pic.height = int(pic.height * factor)
        pic.left = int((SLIDE_W - pic.width) / 2)
    if caption:
        frame = _textbox(slide, MARGIN, pic.top + pic.height + Inches(0.06), BODY_W, caption_h)
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _write(p, caption, 10, MUTED, italic=True)


def _is_group_row(row: list[str]) -> bool:
    return bool(row[0]) and all(not cell for cell in row[1:])


def _table(slide, table: content.Table, top, available):
    n_rows = len(table.rows) + 1
    n_cols = len(table.headers)
    note_h = Inches(0.52) if table.note else Inches(0)
    row_h = min(Inches(0.34), (available - note_h) / n_rows)
    height = int(row_h * n_rows)

    shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, BODY_W, height)
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    size = 12 if n_rows <= 8 and n_cols <= 5 else 10

    # Column widths are proportional rather than fixed: a label column is worth
    # 2.2 numeric ones. A fixed narrow width looks right on the five-column
    # tables and then wraps a header like "Long sample (1,855 days, 3 targets)"
    # onto three lines on the three-column ones.
    aligns = table.alignment()
    weights = [2.2 if a == "l" else 1.0 for a in aligns]
    unit = BODY_W / sum(weights)
    for col, weight in zip(tbl.columns, weights):
        col.width = int(unit * weight)

    def fill(cell, text, *, header=False, group=False, align="l"):
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BAND if (header or group) else WHITE
        frame = cell.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
        colour = ACCENT if header else (MUTED if group else INK)
        _write(p, text, size, colour, bold=header, italic=group)

    for j, head in enumerate(table.headers):
        fill(tbl.cell(0, j), head, header=True, align=aligns[j])

    for i, row in enumerate(table.rows):
        group = _is_group_row(row)
        if group:
            cell = tbl.cell(i + 1, 0)
            cell.merge(tbl.cell(i + 1, n_cols - 1))
            fill(cell, row[0], group=True, align="l")
            continue
        for j, value in enumerate(row):
            fill(tbl.cell(i + 1, j), value, align=aligns[j])

    if table.note:
        frame = _textbox(slide, MARGIN, top + height + Inches(0.08), BODY_W, note_h)
        p = frame.paragraphs[0]
        p.line_spacing = 1.0
        _write(p, "Notes: " + table.note, 9, MUTED, italic=True)

    return int(height + note_h + Inches(0.16))


# ---------------------------------------------------------------------------
# Slide kinds
# ---------------------------------------------------------------------------

def _title_slide(prs, item: content.Slide):
    slide = _blank(prs)
    band = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    frame = _textbox(slide, MARGIN, Inches(2.15), BODY_W, Inches(1.9))
    p = frame.paragraphs[0]
    _write(p, item.title, 36, INK, bold=True)

    frame = _textbox(slide, MARGIN, Inches(3.95), BODY_W, Inches(0.55))
    p = frame.paragraphs[0]
    _write(p, item.subtitle, 20, ACCENT)

    _rule(slide, Inches(4.75), width=Inches(3.2))

    frame = _textbox(slide, MARGIN, Inches(5.05), BODY_W, Inches(0.9))
    p = frame.paragraphs[0]
    _write(p, item.author, 17, INK)
    p = frame.add_paragraph()
    p.space_before = Pt(4)
    _write(p, item.date, 13, MUTED)
    return slide


def _content_slide(prs, item: content.Slide, number: int, total: int):
    slide = _blank(prs)
    _chrome(slide, item, number, total)
    top = BODY_Y
    available = BODY_BOTTOM - BODY_Y

    if item.table is not None:
        # A table is the exhibit: it gets the top of the body and the bullets
        # take what is left, at whatever size that leaves them.
        bullets_budget = Inches(0.62) * len(item.bullets) if item.bullets else Inches(0)
        used = _table(slide, item.table, top, available - bullets_budget)
        top += used
        if item.bullets:
            size = 12 if len(item.bullets) > 2 else 13
            _bullets(slide, item.bullets, top, BODY_BOTTOM - top, size=size)
        return slide

    if item.figure:
        text_h = Inches(0.52) * len(item.bullets) + Inches(0.35)
        _bullets(slide, item.bullets, top, text_h, size=13)
        fig_top = top + text_h
        _figure(slide, item.figure, item.caption, fig_top, BODY_BOTTOM - fig_top)
        return slide

    _bullets(slide, item.bullets, top, available)
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(content.MAIN)
    for i, item in enumerate(content.DECK):
        if item.kind == "title":
            slide = _title_slide(prs, item)
        else:
            slide = _content_slide(prs, item, i + 1, total)
        if item.notes:
            slide.notes_slide.notes_text_frame.text = item.notes

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path} ({len(content.DECK)} slides: "
          f"{len(content.MAIN)} main + {len(content.BACKUP_SLIDES)} backup)")
